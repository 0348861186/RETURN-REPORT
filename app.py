from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN


# ============================================================
# 1. CREATE PRESENTATION - 16:9
# ============================================================

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]


# ============================================================
# 2. COLOR PALETTE
# ============================================================

COLOR_BG = RGBColor(245, 247, 250)
COLOR_BLUE = RGBColor(37, 99, 235)
COLOR_GREEN = RGBColor(16, 185, 129)
COLOR_ORANGE = RGBColor(245, 158, 11)
COLOR_RED = RGBColor(220, 38, 38)
COLOR_DARK = RGBColor(30, 41, 59)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_GRAY = RGBColor(100, 116, 139)


# ============================================================
# 3. HELPER FUNCTIONS
# ============================================================

def set_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG


def add_title_banner(slide, title_en, title_vi, color):
    banner = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(0.4),
        Inches(12.333),
        Inches(0.8)
    )

    banner.fill.solid()
    banner.fill.fore_color.rgb = color
    banner.line.color.rgb = color

    tf = banner.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = title_en
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = title_vi
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_WHITE
    p2.alignment = PP_ALIGN.CENTER


def add_metric_card(
    slide,
    x,
    y,
    width,
    number,
    title_en,
    title_vi,
    color
):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(1.45)
    )

    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = color

    tf = card.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = title_en
    p2.font.size = Pt(11)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_DARK
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = title_vi
    p3.font.size = Pt(9)
    p3.font.color.rgb = COLOR_GRAY
    p3.alignment = PP_ALIGN.CENTER


# ============================================================
# SLIDE 1 - COVER
# ============================================================

slide = prs.slides.add_slide(blank_layout)
set_background(slide)

title = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1.0),
    Inches(1.7),
    Inches(11.3),
    Inches(2.2)
)

title.fill.solid()
title.fill.fore_color.rgb = COLOR_BLUE
title.line.color.rgb = COLOR_BLUE

tf = title.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "WEEKLY CUSTOMER COMPLAINT REPORT"
p.font.size = Pt(30)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "BÁO CÁO KHIẾU NẠI KHÁCH HÀNG HÀNG TUẦN"
p2.font.size = Pt(18)
p2.font.color.rgb = COLOR_WHITE
p2.alignment = PP_ALIGN.CENTER

p3 = tf.add_paragraph()
p3.text = "Week 32 / 2026"
p3.font.size = Pt(14)
p3.font.color.rgb = COLOR_WHITE
p3.alignment = PP_ALIGN.CENTER


# ============================================================
# SLIDE 2 - EXECUTIVE SUMMARY
# ============================================================

slide = prs.slides.add_slide(blank_layout)
set_background(slide)

add_title_banner(
    slide,
    "EXECUTIVE SUMMARY",
    "TÓM TẮT TÌNH HÌNH KHIẾU NẠI",
    COLOR_GREEN
)

add_metric_card(
    slide, 0.5, 1.5, 3.9,
    "3",
    "TOTAL COMPLAINTS",
    "Tổng số khiếu nại",
    COLOR_BLUE
)

add_metric_card(
    slide, 4.7, 1.5, 3.9,
    "2",
    "OPEN CASES",
    "Khiếu nại đang xử lý",
    COLOR_ORANGE
)

add_metric_card(
    slide, 8.9, 1.5, 3.9,
    "1",
    "CLOSED CASES",
    "Khiếu nại đã đóng",
    COLOR_GREEN
)

# Key message
box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5),
    Inches(3.4),
    Inches(12.3),
    Inches(2.2)
)

box.fill.solid()
box.fill.fore_color.rgb = COLOR_WHITE
box.line.color.rgb = COLOR_GRAY

tf = box.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "KEY MESSAGE / THÔNG TIN CHÍNH"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_BLUE

p2 = tf.add_paragraph()
p2.text = (
    "3 customer complaints were received during the week. "
    "The main issue was shortage-related."
)
p2.font.size = Pt(14)
p2.font.color.rgb = COLOR_DARK

p3 = tf.add_paragraph()
p3.text = (
    "Trong tuần phát sinh 3 khiếu nại khách hàng. "
    "Vấn đề chính liên quan đến thiếu số lượng."
)
p3.font.size = Pt(12)
p3.font.color.rgb = COLOR_DARK


# ============================================================
# SLIDE 3 - OVERVIEW & KEY METRICS
# ============================================================

slide = prs.slides.add_slide(blank_layout)
set_background(slide)

add_title_banner(
    slide,
    "OVERVIEW & KEY METRICS",
    "TỔNG QUAN & CHỈ SỐ KPI HÀNG TUẦN",
    COLOR_BLUE
)

add_metric_card(
    slide, 0.5, 1.4, 3.9,
    "3",
    "TOTAL COMPLAINTS RECEIVED",
    "Tổng số khiếu nại phát sinh",
    COLOR_BLUE
)

add_metric_card(
    slide, 4.7, 1.4, 3.9,
    "1",
    "CRITICAL CASES",
    "Khiếu nại nghiêm trọng",
    COLOR_RED
)

add_metric_card(
    slide, 8.9, 1.4, 3.9,
    "33%",
    "CLOSURE RATE",
    "Tỷ lệ hoàn thành",
    COLOR_GREEN
)


# Complaint table
table_shape = slide.shapes.add_table(
    4,
    4,
    Inches(0.5),
    Inches(3.25),
    Inches(12.3),
    Inches(2.7)
)

table = table_shape.table

headers = [
    "Complaint Code",
    "Facility",
    "Customer",
    "Status"
]

data = [
    ["CCR-2026-001", "Plant A", "Customer A", "OPEN"],
    ["CCR-2026-002", "Plant B", "Customer B", "CLOSED"],
    ["CCR-2026-003", "Plant A", "Customer C", "OPEN"]
]

# Header
for i, name in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = name

    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_BLUE

    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

# Data
for row_idx, row_data in enumerate(data, start=1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value

        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_WHITE

        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_DARK
            p.alignment = PP_ALIGN.CENTER


# ============================================================
# SLIDE 4 - PARETO ANALYSIS
# ============================================================

slide = prs.slides.add_slide(blank_layout)
set_background(slide)

add_title_banner(
    slide,
    "PARETO 80/20 ANALYSIS",
    "PHÂN TÍCH PARETO - NỔI BẬT LỖI TRỌNG YẾU",
    COLOR_RED
)

# Chart
chart_data = CategoryChartData()

chart_data.categories = [
    "Shortage",
    "Others"
]

chart_data.add_series(
    "Defects",
    (66.7, 33.3)
)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.DOUGHNUT,
    Inches(6.8),
    Inches(1.5),
    Inches(5.5),
    Inches(4.5),
    chart_data
).chart

chart.has_legend = True
chart.legend.include_in_layout = False

chart.has_title = True
chart.chart_title.text_frame.text = (
    "DEFECT CATEGORY / PHÂN LOẠI LỖI"
)

# Main finding
finding = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.6),
    Inches(1.6),
    Inches(5.6),
    Inches(3.8)
)

finding.fill.solid()
finding.fill.fore_color.rgb = COLOR_WHITE
finding.line.color.rgb = COLOR_RED

tf = finding.text_frame
tf.clear()

p = tf.paragraphs[0]
p.text = "KEY FINDING"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLOR_RED

p2 = tf.add_paragraph()
p2.text = "Shortage = 66.7%"
p2.font.size = Pt(26)
p2.font.bold = True
p2.font.color.rgb = COLOR_RED

p3 = tf.add_paragraph()
p3.text = (
    "Shortage-related complaints represent "
    "the largest defect category."
)
p3.font.size = Pt(13)
p3.font.color.rgb = COLOR_DARK

p4 = tf.add_paragraph()
p4.text = (
    "Khiếu nại liên quan đến thiếu số lượng "
    "chiếm tỷ trọng lớn nhất."
)
p4.font.size = Pt(12)
p4.font.color.rgb = COLOR_DARK


# ============================================================
# SLIDE 5 - ACTION PLAN
# ============================================================

slide = prs.slides.add_slide(blank_layout)
set_background(slide)

add_title_banner(
    slide,
    "CORRECTIVE ACTION & FOLLOW-UP",
    "HÀNH ĐỘNG KHẮC PHỤC & THEO DÕI",
    COLOR_ORANGE
)

actions = [
    [
        "1",
        "Shortage investigation",
        "Điều tra nguyên nhân thiếu số lượng",
        "Production / QA",
        "OPEN"
    ],
    [
        "2",
        "Verify packing process",
        "Xác nhận lại quy trình đóng gói",
        "Production",
        "OPEN"
    ],
    [
        "3",
        "Customer feedback",
        "Phản hồi và xác nhận với khách hàng",
        "QA",
        "PENDING"
    ],
    [
        "4",
        "Effectiveness verification",
        "Xác nhận hiệu lực hành động",
        "QA",
        "PENDING"
    ]
]

table_shape = slide.shapes.add_table(
    len(actions) + 1,
    5,
    Inches(0.5),
    Inches(1.5),
    Inches(12.3),
    Inches(4.8)
)

table = table_shape.table

headers = [
    "No.",
    "Action",
    "Hành động",
    "Responsible",
    "Status"
]

for col, header in enumerate(headers):
    cell = table.cell(0, col)
    cell.text = header

    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_ORANGE

    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

for row, action in enumerate(actions, start=1):
    for col, value in enumerate(action):
        cell = table.cell(row, col)
        cell.text = value

        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_WHITE

        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_DARK
            p.alignment = PP_ALIGN.CENTER


# ============================================================
# 6. SAVE FILE
# ============================================================

output_file = "Weekly_CCR_Report.pptx"

prs.save(output_file)

print(f"Đã xuất file PowerPoint thành công: {output_file}")
